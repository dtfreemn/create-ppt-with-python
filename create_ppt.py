# for path manipulation
from os import getcwd

# for google api calls
from apiclient import discovery
from httplib2 import Http
from credentials import get_credentials # relative import of function to verify credentials

# for making ppts
from pptx import Presentation
from pptx.util import Inches

# for sending email
from smtplib import SMTP
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email.encoders import encode_base64

# relative import of private variables
from secrets import FROM_EMAIL, FROM_PASSWORD, TO_EMAIL, SHEET_ID

def convert_file_to_attachment(file):
	file_name = file.split('/')[-1] # extract file name from absolute path
	file_object = open(file, 'rb')
	attachment = MIMEBase('application', 'octet-stream')
	attachment.set_payload(file_object.read())
	encode_base64(attachment)
	attachment.add_header('Content-Disposition', "attachment; filename= %s" % file_name)
	return attachment

def make_email(files_to_attach):
	email = MIMEMultipart()
	email_body = 'See attached.'
	email.attach(MIMEText(email_body, 'plain'))

	email['From'] = FROM_EMAIL
	email['To'] = TO_EMAIL
	email['Subject'] = 'This Week\'s Community Meeting Powerpoints'

	for file in files_to_attach:
		attachment = convert_file_to_attachment(file)
		email.attach(attachment)

	packaged_email = email.as_string()

	return packaged_email

def send_email(email):
	server = SMTP('smtp.gmail.com', 587)
	server.starttls()
	server.login(FROM_EMAIL, FROM_PASSWORD)
	server.sendmail(FROM_EMAIL, TO_EMAIL, email)
	server.quit()

def proccess_tabs(tabs, service):
	parsed_tabs = []

	for tab in tabs:
		tab_name = tab.get('properties', {}).get('title')
		range_name = tab_name + '!A2:C'
		tab_data = service.spreadsheets().values().get(spreadsheetId=SHEET_ID, range=range_name).execute()
		student_data = tab_data.get('values', [])
		file_path = getcwd() + '/PowerPoints/' + tab_name + '.pptx'
		tab_object = {'file': file_path, 'student_data': student_data}
		parsed_tabs.append(tab_object)

	return parsed_tabs

def make_all_ppts(parsed_tabs):
	for tab in parsed_tabs:
		file_path = tab['file']
		student_data = tab['student_data']
		make_ppt(file_path, student_data)

def make_ppt(file_path, student_data):
	presentation = Presentation()
	blank_slide_layout = presentation.slide_layouts[6]
	slide = presentation.slides.add_slide(blank_slide_layout)

	r = 0
	placeholder = 0
	left = 0.5
	top = 1
	width = 1
	height = 1

	for row in student_data:
		if r > 0:
			if placeholder == 34:
				slide = presentation.slides.add_slide(blank_slide_layout)
				placeholder = 0
				left = 0.5
				top = 1
			if placeholder == 12 or placeholder == 24:
				left = left + 3
				top = 1
			first_name = row[0]
			last_name = row[1]
			pwa = row[2]
			textBox = slide.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(height))
			textBox.text_frame.text = first_name + ' ' + last_name + ' ' + pwa
			placeholder = placeholder+1
			top = top+0.5
		r = r+1

	presentation.save(file_path)

def isolate_file_names_for_attachments(parsed_tabs):
	file_names = []
	
	for tab in parsed_tabs:
		file_names.append(tab['file'])

	return file_names

def main():
  credentials = get_credentials()
  http = credentials.authorize(Http())
  discoveryUrl = ('https://sheets.googleapis.com/$discovery/rest?''version=v4')
  service = discovery.build('sheets', 'v4', http=http, discoveryServiceUrl=discoveryUrl)
  tabs = service.spreadsheets().get(spreadsheetId=SHEET_ID).execute().get('sheets', '')
  processed_tabs = proccess_tabs(tabs, service)
  make_all_ppts(processed_tabs)
  files_to_attach = isolate_file_names_for_attachments(processed_tabs)
  email = make_email(files_to_attach)
  send_email(email)


if __name__ == '__main__':
	main()